#!/usr/bin/env python3
"""Build an editable .pptx (with working hyperlinks) from a Marp markdown deck.

Why this exists
---------------
`marp-cli ... -o deck.pptx` rasterises every slide into a PNG
(`ppt/media/Slide-N-image-1.png`). The result has zero text nodes and zero
hyperlinks, so nothing in the deck is clickable or editable in PowerPoint.
This script emits native text frames and tables via python-pptx instead, so
`[label](url)` becomes a real hyperlink relationship.

Use this — not marp-cli — whenever the deck's links need to work.

Usage
-----
    python3 tools/build-slide-pptx.py <input.md> <output.pptx>

    # example
    python3 tools/build-slide-pptx.py \\
        1_MeetingMinutes/0821_slide.md 1_MeetingMinutes/0821_slide.pptx

Requires `python-pptx` (pip install python-pptx).

Supported markdown
------------------
YAML front matter (`header:` / `footer:` are read), `---` slide separators,
`<!-- _class: lead -->` for centred cover/divider slides, `#` titles, tables,
bullet and numbered lists, ```code fences```, `>` blockquote notes, and the
inline spans `**bold**`, `` `code` `` and `[label](url)`.

Links must be absolute URLs: PowerPoint has no base directory, so relative
paths such as `../FOO.md` will not resolve for the reader.
"""
import re
import sys
from pathlib import Path
from urllib.parse import urlsplit, urlunsplit, quote

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

if len(sys.argv) != 3:
    sys.exit(f"usage: {Path(sys.argv[0]).name} <input.md> <output.pptx>")

SRC = Path(sys.argv[1])
DST = Path(sys.argv[2])

if not SRC.is_file():
    sys.exit(f"input not found: {SRC}")

FONT = "Microsoft JhengHei"
MONO = "Consolas"
INK = RGBColor(0x22, 0x30, 0x3F)
HEAD = RGBColor(0x1A, 0x2B, 0x3C)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
MUTED = RGBColor(0x8A, 0x96, 0xA3)
LINK = RGBColor(0x1B, 0x5E, 0xA8)
RULE = RGBColor(0xD6, 0xDC, 0xE3)
BAND = RGBColor(0xF7, 0xF9, 0xFA)
HDRFILL = RGBColor(0xEC, 0xEF, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_PER_IN = 914400.0

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.62)
CONTENT_W = SLIDE_W - 2 * MARGIN_L
TEXT_W_IN = CONTENT_W / EMU_PER_IN
FOOT_Y = SLIDE_H - Inches(0.62)
# Lowest y a flow diagram may reach, leaving the footer band clear.
CONTENT_BOTTOM = Inches(6.45)

# The flow block highlights the exit whose text contains this token. Deck-
# specific: change it (or clear it) when reusing this script for another deck.
FLOW_HIGHLIGHT = "Option B"


# ---------------------------------------------------------------- md parsing
def split_slides(text):
    body = text
    if body.startswith("---"):
        body = body[body.index("\n---", 3) + 4:]
    return [p.strip("\n") for p in re.split(r"\n---[ \t]*\n", body) if p.strip()]


def parse(raw):
    lines = raw.split("\n")
    blocks, i, lead, as_flow = [], 0, False, False
    while i < len(lines):
        s = lines[i].strip()
        if not s:
            i += 1
            continue
        if s.startswith("<!--"):
            lead = lead or ("lead" in s)
            if re.search(r"\bflow\b", s):
                as_flow = True
            i += 1
            continue
        if s.startswith("# "):
            blocks.append(("title", s[2:].strip()))
            i += 1
            continue
        if s.startswith("```"):
            i += 1
            code = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            blocks.append(("code", code))
            continue
        if s.startswith("|"):
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                r = lines[i].strip()
                i += 1
                if re.match(r"^\|[\s:|-]+\|$", r):
                    continue
                rows.append([c.strip() for c in r.strip("|").split("|")])
            if rows:
                blocks.append(("flow" if as_flow else "table", rows))
                as_flow = False
            continue
        if s.startswith(">"):
            q = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                q.append(lines[i].strip().lstrip(">").strip())
                i += 1
            blocks.append(("note", [x for x in q if x]))
            continue
        if LIST_ITEM.match(s):
            kind = "numbers" if s[0].isdigit() else "bullets"
            it = []
            while i < len(lines):
                cur, st = lines[i], lines[i].strip()
                # Indentation wins over the item pattern: an indented "- x" is
                # a nested sub-item, not a new top-level entry.
                if cur.startswith("  ") and it:
                    it[-1] += "\n" + st
                    i += 1
                    continue
                m = LIST_ITEM.match(st)
                if m:
                    it.append(m.group(2).strip())
                    i += 1
                    continue
                if not st:
                    # A blank line only ends the list when what follows is
                    # neither another item nor an indented continuation.
                    j = i + 1
                    while j < len(lines) and not lines[j].strip():
                        j += 1
                    if j < len(lines) and (LIST_ITEM.match(lines[j].strip())
                                           or lines[j].startswith("  ")):
                        i = j
                        continue
                    break
                break
            blocks.append((kind, it))
            continue
        if re.fullmatch(r"\*\*.+\*\*", s):
            blocks.append(("label", s.strip("*")))
            i += 1
            continue
        blocks.append(("para", s))
        i += 1
    return blocks, lead


# ------------------------------------------------------------- run rendering
LIST_ITEM = re.compile(r"^([-*]|\d+\.)\s+(.*)$")

TOKEN = re.compile(r"(\[[^\]]+\]\([^)]*\)|\*\*.+?\*\*|`[^`]+`)")


def encode_url(url):
    """Percent-encode non-ASCII path/fragment.

    OOXML relationship targets must be valid URIs. A raw UTF-8 fragment such as
    `#一頁結論` leaves the link dead in PowerPoint, so encode path and fragment
    while leaving already-escaped sequences and normal path punctuation intact.
    """
    parts = urlsplit(url)
    if not parts.scheme:
        return url
    return urlunsplit((
        parts.scheme,
        parts.netloc,
        quote(parts.path, safe="/-._~+()!*'"),
        quote(parts.query, safe="=&-._~+"),
        quote(parts.fragment, safe="-._~"),
    ))


def emit(para, text, size, color=INK, bold=False):
    """Render inline markdown spans as runs; real hyperlinks for [x](url)."""
    for piece in TOKEN.split(text):
        if not piece:
            continue
        m = re.fullmatch(r"\[([^\]]+)\]\(([^)]*)\)", piece)
        if m:
            label, url = m.group(1), m.group(2)
            label = label.replace("**", "").replace("`", "")
            r = para.add_run()
            r.text = label
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            if url:
                r.hyperlink.address = encode_url(url)
            r.font.color.rgb = LINK
            r.font.underline = True
            continue
        r = para.add_run()
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if piece.startswith("**") and piece.endswith("**"):
            r.text = piece[2:-2]
            r.font.bold = True
            r.font.color.rgb = ACCENT
        elif piece.startswith("`") and piece.endswith("`"):
            r.text = piece[1:-1]
            r.font.name = MONO
        else:
            r.text = piece


def box(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def fill_cell(cell, text, size, bold, align, bg):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.margin_left = cell.margin_right = Inches(0.07)
    cell.margin_top = cell.margin_bottom = Inches(0.035)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    emit(p, text, size, HEAD if bold else INK, bold=bold)


def _vis(s):
    """Text as it will render: markdown markers stripped, link labels kept."""
    s = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", s)
    return s.replace("**", "").replace("`", "")


def _width_pt(s, size):
    """Approximate rendered width in points (CJK ~1em, latin ~0.55em)."""
    return sum(size if ord(ch) > 0x2E80 else size * 0.55 for ch in s)


CELL_PAD_IN = 0.07  # matches fill_cell left/right margin
ROW_LINE_FACTOR = 1.45  # table rows: glyph box plus in-cell leading
TABLE_SZ_DELTA = 5      # table type runs this many pt below body type
MIN_FLOW_PT = 9         # smallest type a flow diagram will shrink to
MIN_BAND_EMU = int(Inches(0.34))  # floor for a table row / flow band
NOTE_RESERVE_EMU = int(Inches(0.55))  # space a flow block leaves for a following note


def line_h(size, space_after_pt=6.0):
    """Rendered height of one text line, in EMU (glyph box + paragraph spacing)."""
    return int(((size * 1.25 + space_after_pt) / 72.0) * EMU_PER_IN)


def wrap_lines(text, width_in, size):
    """Wrapped line count for `text` in a box `width_in` inches wide.

    Must measure the *visible* text: a markdown link carries a long absolute
    URL that never renders, so counting raw characters wildly overestimates
    the height and pushes later blocks off the slide.
    """
    usable_pt = max(width_in * 72, 1.0)
    need = _width_pt(_vis(text), size)
    return max(1, int(need / usable_pt) + (1 if need % usable_pt else 0))


def layout_table(rows, total_emu, size):
    """Column widths plus per-row heights that account for wrapped lines.

    A fixed row height makes tall (wrapped) rows overflow their allotted space,
    which pushes the rendered table past the geometry we reserved and lets the
    next block land on top of it. Estimating the wrap per row keeps the flow
    honest.
    """
    n = max(len(r) for r in rows)
    norm = [r + [""] * (n - len(r)) for r in rows]

    # Dampen the weight so one very long cell cannot starve the other columns.
    weights = []
    for c in range(n):
        cells = [_width_pt(_vis(r[c]), size) for r in norm]
        weights.append(max(0.6 * max(cells) + 0.4 * (sum(cells) / len(cells)),
                           size * 3.0))
    tot = sum(weights)
    cw = [int(total_emu * w / tot) for w in weights]

    # A column must fit its longest unbreakable token, otherwise latin words
    # such as "CockroachDB" get split mid-word. CJK wraps anywhere, so only
    # latin/digit runs matter here. Cap the floor so one long token cannot
    # swallow the table.
    pad_x_emu = int(Inches(CELL_PAD_IN) * 2)
    cap = int(total_emu * 0.34)
    floors = []
    for c in range(n):
        longest = 0.0
        for r in norm:
            for tok in re.findall(r"[0-9A-Za-z][0-9A-Za-z._/＋+-]*", _vis(r[c])):
                longest = max(longest, _width_pt(tok, size))
        floors.append(min(int((longest / 72.0) * EMU_PER_IN) + pad_x_emu, cap))

    # Raise any column below its floor, then reclaim from those with slack.
    for _ in range(4):
        deficit = sum(max(0, floors[c] - cw[c]) for c in range(n))
        if deficit <= 0:
            break
        slack = [max(0, cw[c] - floors[c]) for c in range(n)]
        total_slack = sum(slack)
        if total_slack <= 0:
            break
        for c in range(n):
            if cw[c] < floors[c]:
                cw[c] = floors[c]
            elif slack[c]:
                cw[c] -= int(deficit * slack[c] / total_slack)
    # Keep the row total exactly on the table width.
    drift = total_emu - sum(cw)
    cw[max(range(n), key=lambda c: cw[c])] += drift

    # Row line height carries the cell's own leading, so it is deliberately
    # looser than line_h()'s paragraph spacing.
    line_pt = size * ROW_LINE_FACTOR
    pad_y_emu = int(Inches(0.035) * 2)
    heights = []
    for r in norm:
        lines = 1
        for c in range(n):
            usable_in = cw[c] / EMU_PER_IN - CELL_PAD_IN * 2
            if usable_in * 72 <= 1:
                continue
            lines = max(lines, wrap_lines(r[c], usable_in, size))
        h = int((lines * line_pt / 72) * EMU_PER_IN) + pad_y_emu
        heights.append(max(h, MIN_BAND_EMU))
    return cw, n, norm, heights



FLOW_FILL = RGBColor(0xF2, 0xF5, 0xF8)
FLOW_LINE = RGBColor(0xB9, 0xC4, 0xCF)
FLOW_PICK = RGBColor(0xFD, 0xF1, 0xEF)


def _flow_box(slide, x, y, w, h, text, size, bold=False, accent=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = FLOW_PICK if accent else FLOW_FILL
    sp.line.color.rgb = ACCENT if accent else FLOW_LINE
    sp.line.width = Pt(1.5 if accent else 0.75)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    emit(p, text, size, bold=bold)
    return sp


def _flow_arrow(slide, x, y, w, h):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(x), int(y), int(w), int(h))
    a.fill.solid()
    a.fill.fore_color.rgb = FLOW_LINE
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def _draw_flow(slide, rows, y0, size, avail):
    """Render a 3-column markdown table as a branching flow diagram.

    A table forces linear reading, but a decision with shared exits reads as a
    flat list. Grouping the first column across its branches makes the actual
    structure visible.

    Band heights are per row: a uniform height clips whichever branch needs two
    lines, so each row gets what its longest cell actually needs and the whole
    stack is then scaled to the space available.
    """
    body = rows[1:] if len(rows) > 1 else rows
    n = len(body)
    gap = int(Inches(0.10))
    arrow_w = int(Inches(0.26))
    col = [int(CONTENT_W * 0.29), int(CONTENT_W * 0.24), 0]
    x0 = int(MARGIN_L)
    x1 = x0 + col[0] + arrow_w + gap
    x2 = x1 + col[1] + arrow_w + gap
    col[2] = int(CONTENT_W) - (x2 - x0)
    pad = int(Inches(0.09))

    def plan(sz):
        lh = line_h(sz, 0)
        hs = []
        for r in body:
            need = 1
            for c, txt in enumerate((r + ["", ""])[:3]):
                cw_in = col[c] / EMU_PER_IN - CELL_PAD_IN * 2
                need = max(need, wrap_lines(txt, cw_in, sz if c != 1 else sz - 1))
            hs.append(max(need * lh + pad, int(Inches(0.42))))
        return hs

    smallest = max(MIN_FLOW_PT, size - 3)
    for sz in range(size, smallest - 1, -1):
        heights = plan(sz)
        if sum(heights) + gap * (n - 1) <= avail:
            break
    size = sz
    if sum(heights) + gap * (n - 1) > avail:
        # Still too tall at the smallest type: tighten the gap, then scale.
        gap = int(Inches(0.05))
        total = sum(heights) + gap * (n - 1)
        if total > avail:
            k = (avail - gap * (n - 1)) / float(sum(heights))
            heights = [max(int(h * k), MIN_BAND_EMU) for h in heights]

    tops, acc = [], y0
    for h in heights:
        tops.append(acc)
        acc += h + gap

    groups, cur = [], None
    for idx, r in enumerate(body):
        if _vis(r[0]).strip():
            cur = [idx, idx, r[0]]
            groups.append(cur)
        elif cur:
            cur[1] = idx

    for a, b, label in groups:
        gh = tops[b] + heights[b] - tops[a]
        _flow_box(slide, x0, tops[a], col[0], gh, label, size, bold=True)

    for idx, r in enumerate(body):
        by, bh = tops[idx], heights[idx]
        cond = r[1] if len(r) > 1 else ""
        exit_ = r[2] if len(r) > 2 else ""
        pick = bool(FLOW_HIGHLIGHT) and FLOW_HIGHLIGHT in _vis(exit_)
        mid = by + bh // 2 - int(Inches(0.055))
        if _vis(cond).strip() in ("", "—", "-"):
            _flow_arrow(slide, x0 + col[0] + gap // 2, mid,
                        col[1] + arrow_w + gap, int(Inches(0.11)))
        else:
            _flow_arrow(slide, x0 + col[0] + gap // 2, mid, arrow_w, int(Inches(0.11)))
            _flow_box(slide, x1, by, col[1], bh, cond, size - 1)
            _flow_arrow(slide, x1 + col[1] + gap // 2, mid, arrow_w, int(Inches(0.11)))
        _flow_box(slide, x2, by, col[2], bh, exit_, size, accent=pick)

    return sum(heights) + gap * (n - 1) + int(Inches(0.20))


# ---------------------------------------------------------------- build deck
def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]
    md = SRC.read_text()
    header = re.search(r"^header:\s*'(.*)'$", md, re.M)
    footer = re.search(r"^footer:\s*'(.*)'$", md, re.M)
    HDR = header.group(1) if header else ""
    FTR = footer.group(1) if footer else ""

    for idx, raw in enumerate(split_slides(md), start=1):
        blocks, lead = parse(raw)
        s = prs.slides.add_slide(blank)

        if idx > 1 and HDR:
            emit(box(s, MARGIN_L, Inches(0.26), CONTENT_W, Inches(0.3)).paragraphs[0],
                 HDR, 11, MUTED)
        if FTR:
            emit(box(s, MARGIN_L, FOOT_Y, Inches(4), Inches(0.3)).paragraphs[0], FTR, 11, MUTED)
        tf = box(s, SLIDE_W - MARGIN_L - Inches(1), FOOT_Y, Inches(1), Inches(0.3))
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        emit(tf.paragraphs[0], str(idx), 12, MUTED)

        weight = 0
        for k, pl in blocks:
            if k == "table":
                weight += 2 + len(pl)
            elif k in ("bullets", "numbers"):
                weight += sum(1 + wrap_lines(x, TEXT_W_IN - 0.3, 18) for x in pl)
            elif k in ("code", "note"):
                weight += len(pl)
            elif k in ("para", "label"):
                weight += wrap_lines(pl, TEXT_W_IN, 18)
        body_sz = 17 if weight >= 18 else (18 if weight >= 13 else 19)
        tbl_sz = body_sz - TABLE_SZ_DELTA

        if lead:
            y = Inches(2.7)
            for k, pl in blocks:
                if k == "title":
                    tf = box(s, MARGIN_L, y, CONTENT_W, Inches(1.1))
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    emit(tf.paragraphs[0], pl, 30, HEAD, bold=True)
                    y += Inches(1.25)
                elif k in ("para", "label"):
                    tf = box(s, MARGIN_L, y, CONTENT_W, Inches(0.42))
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    emit(tf.paragraphs[0], pl, 16, MUTED)
                    y += Inches(0.46)
            continue

        y = Inches(0.95)
        for bi, (k, pl) in enumerate(blocks):
            if k == "title":
                emit(box(s, MARGIN_L, y, CONTENT_W, Inches(0.62)).paragraphs[0],
                     pl, 25, HEAD, bold=True)
                y += Inches(0.92)

            elif k == "label":
                emit(box(s, MARGIN_L, y, CONTENT_W, Inches(0.34)).paragraphs[0],
                     pl, body_sz, ACCENT, bold=True)
                y += Inches(0.44)

            elif k == "para":
                ln = wrap_lines(pl, TEXT_W_IN, body_sz)
                hh = line_h(body_sz, 0) * ln
                emit(box(s, MARGIN_L, y, CONTENT_W, hh).paragraphs[0], pl, body_sz)
                y += hh + Inches(0.14)

            elif k in ("bullets", "numbers"):
                h = Inches(0.0)
                tf = box(s, MARGIN_L, y, CONTENT_W, Inches(0.4) * len(pl))
                for n, item in enumerate(pl):
                    parts = item.split("\n")
                    p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
                    p.space_after = Pt(6)
                    r = p.add_run()
                    r.text = f"{n+1}. " if k == "numbers" else "•  "
                    r.font.name = FONT
                    r.font.size = Pt(body_sz)
                    r.font.color.rgb = MUTED
                    emit(p, parts[0], body_sz)
                    h += line_h(body_sz) * wrap_lines(parts[0], TEXT_W_IN - 0.3, body_sz)
                    for extra in parts[1:]:
                        p2 = tf.add_paragraph()
                        p2.space_after = Pt(6)
                        rr = p2.add_run()
                        rr.text = "     "
                        rr.font.size = Pt(body_sz)
                        emit(p2, extra, body_sz - 1, MUTED)
                        h += line_h(body_sz - 1) * wrap_lines(extra, TEXT_W_IN - 0.5, body_sz - 1)
                y += h + Inches(0.12)

            elif k == "code":
                h = Inches(0.26) * len(pl) + Inches(0.24)
                bg = s.shapes.add_shape(1, MARGIN_L, y, int(CONTENT_W), int(h))
                bg.fill.solid()
                bg.fill.fore_color.rgb = BAND
                bg.line.color.rgb = RULE
                bg.shadow.inherit = False
                tf = box(s, MARGIN_L + Inches(0.14), y + Inches(0.12),
                         CONTENT_W - Inches(0.28), h - Inches(0.24))
                for n, ln in enumerate(pl):
                    p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
                    p.space_after = Pt(0)
                    r = p.add_run()
                    r.text = ln if ln.strip() else " "
                    r.font.name = MONO
                    r.font.size = Pt(body_sz - 4)
                    r.font.color.rgb = INK
                y += h + Inches(0.16)

            elif k == "note":
                nsz = body_sz - 3
                nl = sum(wrap_lines(q, TEXT_W_IN - 0.2, nsz) for q in pl)
                nh = line_h(nsz, 2) * nl
                tf = box(s, MARGIN_L + Inches(0.14), y + Inches(0.04),
                         CONTENT_W - Inches(0.14), nh)
                for n, q in enumerate(pl):
                    p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
                    emit(p, q, nsz, MUTED)
                bar = s.shapes.add_shape(1, MARGIN_L, y + Inches(0.04),
                                         Emu(26000), nh)
                bar.fill.solid()
                bar.fill.fore_color.rgb = RULE
                bar.line.fill.background()
                bar.shadow.inherit = False
                y += nh + Inches(0.18)

            elif k == "flow":
                # Reserve room for a trailing note. Use the loop index, not a
                # value search: two identical blocks on one slide would
                # otherwise resolve to the first one's position.
                tail = any(kk == "note" for kk, _ in blocks[bi + 1:])
                avail = int(CONTENT_BOTTOM) - y - (NOTE_RESERVE_EMU if tail else 0)
                y += _draw_flow(s, pl, y, tbl_sz, avail)

            elif k == "table":
                cw, ncol, norm, rhs = layout_table(pl, int(CONTENT_W), tbl_sz)
                nrow = len(norm)
                tbl = s.shapes.add_table(nrow, ncol, MARGIN_L, y,
                                         int(CONTENT_W), sum(rhs)).table
                tbl.first_row = True
                for c in range(ncol):
                    tbl.columns[c].width = cw[c]
                # Alignment is decided per column, not per cell: judging each
                # cell on its own length leaves a column visually ragged
                # (short values centred, longer ones flush left).
                body = norm[1:] or norm
                aligns = [
                    PP_ALIGN.CENTER
                    if max(len(_vis(r[c])) for r in body) <= 14
                    else PP_ALIGN.LEFT
                    for c in range(ncol)
                ]
                for r in range(nrow):
                    tbl.rows[r].height = rhs[r]
                    for c in range(ncol):
                        hd = (r == 0)
                        al = PP_ALIGN.CENTER if hd else aligns[c]
                        bg = HDRFILL if hd else (BAND if r % 2 == 0 else WHITE)
                        fill_cell(tbl.cell(r, c), norm[r][c], tbl_sz, hd, al, bg)
                y += sum(rhs) + Inches(0.20)

    prs.save(DST)
    print(f"wrote {DST}")


build()
