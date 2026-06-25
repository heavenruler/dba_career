#!/usr/bin/env python3
"""Build 0626-slide-v6.pptx — 17-slide editable deck with dark-navy theme.

Extends v5 (11 slides) with 6 inserted pages aligned to the same theme:
  - Slide 7/8/9: YBDB / CRDB / TiDB architecture (extracted image + native bullets)
  - Slide 10: commit path × cross-region latency tax × consistency
  - Slide 11: range vs hash sharding
  - Slide 14: P-A / P-B placement intro

Inputs: 1_MeetingMinutes/assets/arch-{ybdb,crdb,tidb}-1.{jpg,png}
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---- palette ------------------------------------------------------------
BG_NAVY     = RGBColor(0x1E, 0x25, 0x38)
ACCENT_RED  = RGBColor(0xE7, 0x4C, 0x3C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLATE       = RGBColor(0xE8, 0xEA, 0xED)
MUTED       = RGBColor(0xA0, 0xA8, 0xB8)
ORANGE      = RGBColor(0xF3, 0x9C, 0x12)
RED_EMPH    = RGBColor(0xB0, 0x3A, 0x14)
TABLE_BG    = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ALT   = RGBColor(0xF1, 0xF3, 0xF6)
TABLE_HEAD  = RGBColor(0xD8, 0xDD, 0xE6)
TABLE_TEXT  = RGBColor(0x11, 0x14, 0x18)
TABLE_BORD  = RGBColor(0x8C, 0x95, 0xA3)

FONT_CJK    = 'PingFang TC'
FONT_LATIN  = 'Inter'

# ---- layout -------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, *, size=22, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font or FONT_CJK
    return box

def add_runs(slide, x, y, w, h, runs, *, size=20, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """runs = list of (text, {bold,color,size,font}) -- all on one paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, opts in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.color.rgb = opts.get('color', SLATE)
        r.font.name = opts.get('font', FONT_CJK)
    return box

def add_bullets(slide, x, y, w, h, items, *, size=18, color=SLATE,
                bullet_color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # XML-level bullet char
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgb = etree.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', f'{(bullet_color or ACCENT_RED).rgb:06X}' if hasattr((bullet_color or ACCENT_RED), 'rgb') else 'E74C3C')
        pPr.set('indent', '-228600')
        pPr.set('marL', '228600')

        if isinstance(item, str):
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = FONT_CJK
        else:
            for txt, opts in item:
                r = p.add_run()
                r.text = txt
                r.font.size = Pt(opts.get('size', size))
                r.font.bold = opts.get('bold', False)
                r.font.color.rgb = opts.get('color', color)
                r.font.name = opts.get('font', FONT_CJK)
    return box

def add_header_footer(slide, page_num, total):
    # top accent bar
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Pt(6), ACCENT_RED)
    # header text
    add_text(slide, Inches(0.4), Inches(0.18), Inches(6), Inches(0.3),
             'PoC 第一階段成果與下一步決策', size=12, color=MUTED)
    # footer
    add_text(slide, Inches(0.4), Inches(7.05), Inches(6), Inches(0.3),
             '2026-06-26 · DBA', size=12, color=MUTED)
    # page number
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3),
             str(page_num), size=12, color=MUTED, align=PP_ALIGN.RIGHT)

def add_title(slide, text, y=Inches(0.55)):
    add_text(slide, Inches(0.6), y, Inches(12), Inches(0.7),
             text, size=30, bold=True, color=WHITE)
    # underline
    add_rect(slide, Inches(0.6), y + Inches(0.7), Inches(12), Pt(2), ACCENT_RED)

def add_subtitle(slide, text, y, *, color=MUTED, size=18):
    """Subtitle with red left bar."""
    bar_x = Inches(0.6)
    bar_w = Pt(3)
    add_rect(slide, bar_x, y, bar_w, Inches(0.4), ACCENT_RED)
    add_text(slide, bar_x + Emu(80000), y, Inches(12), Inches(0.4),
             text, size=size, color=color, anchor=MSO_ANCHOR.MIDDLE)

# ---- table builder ------------------------------------------------------
def style_cell(cell, *, bg=TABLE_BG, color=TABLE_TEXT, bold=False, size=14,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, runs=None):
    """Style a table cell. runs=[(text,opts)...] overrides default text rendering."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.margin_left = cell.margin_right = Inches(0.08)
    cell.margin_top = cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = anchor
    tf = cell.text_frame
    tf.word_wrap = True
    # clear existing text
    for p in list(tf.paragraphs):
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
    p = tf.paragraphs[0]
    p.alignment = align
    if runs is None:
        return  # caller will set text manually via cell.text
    for txt, opts in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', bold)
        r.font.color.rgb = opts.get('color', color)
        r.font.name = opts.get('font', FONT_CJK)

def add_table(slide, x, y, w, h, header, rows, *, col_widths=None,
              header_size=14, body_size=13):
    """header = [str]; rows = [[cell_value]] where cell_value is str or [(text,opts)..]."""
    n_rows = len(rows) + 1
    n_cols = len(header)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    # column widths
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    # header row
    for i, txt in enumerate(header):
        cell = tbl.cell(0, i)
        style_cell(cell, bg=TABLE_HEAD, color=TABLE_TEXT, bold=True,
                   size=header_size, align=PP_ALIGN.CENTER,
                   runs=[(txt, {'bold': True, 'size': header_size})])
    # body rows
    for r_idx, row in enumerate(rows):
        bg = TABLE_BG if r_idx % 2 == 0 else TABLE_ALT
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            if isinstance(val, str):
                style_cell(cell, bg=bg, size=body_size,
                           runs=[(val, {'size': body_size})])
            else:
                # list of (text, opts)
                style_cell(cell, bg=bg, size=body_size, runs=val)
    return tbl

# =========================================================================
# Slide 1 — Title
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_rect(s, Emu(0), Emu(0), SLIDE_W, Pt(6), ACCENT_RED)
add_text(s, Inches(0.4), Inches(0.18), Inches(6), Inches(0.3),
         'PoC 第一階段成果與下一步決策', size=12, color=MUTED)
add_text(s, Inches(0), Inches(3.0), SLIDE_W, Inches(0.9),
         '第一階段 PoC 成果與下一步決策', size=44, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.2), SLIDE_W, Inches(0.5),
         '分散式資料庫架構 × 跨區域驗證 × 決策框架', size=20, color=MUTED,
         align=PP_ALIGN.CENTER)
# divider
add_rect(s, Inches(5), Inches(5.0), Inches(3.3), Pt(2), ACCENT_RED)
add_text(s, Inches(0), Inches(5.2), SLIDE_W, Inches(0.5),
         '2026-06-26    |    DBA', size=18, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3),
         '1', size=12, color=MUTED, align=PP_ALIGN.RIGHT)

# =========================================================================
# Slide 2 — Outline
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 2, 17)
add_title(s, 'Outline')
add_bullets(s, Inches(0.8), Inches(1.7), Inches(10), Inches(5),
            ['專案迄今歷程',
             '第一階段 PoC 測試數據彙整',
             '第二階段 跨區域 / 跨專線執行進度',
             '決策框架說明',
             '後續推進'],
            size=22, color=SLATE, bullet_color=ACCENT_RED)

# =========================================================================
# Slide 3 — 專案歷程 ①
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 3, 17)
add_title(s, '專案歷程 ① 研究定義到跨家框架')
add_subtitle(s, '先定義問題與比較口徑，再建立可重複執行的基礎設施與三資料庫共同工具鏈',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.4),
          header=['時間', '階段', '重大節點', '狀態'],
          rows=[
            ['2026-03-30～04-10',
             [('前期研究', {'bold': True, 'color': ORANGE})],
             '定義分散式 SQL、跨區同鍵寫入、follower read、HA/DR 與九項 survey 評估面向',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-04-21～04-27',
             [('IaC 與第一版測試鏈', {'bold': True, 'color': ORANGE})],
             [('建立多測項部署及獨立壓測流程 (Flow Chain)\n', {}),
              ('gate → prepare → gate-iso → dry-run → run → collect → summary',
               {'color': ORANGE, 'size': 11})],
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-04-28～05-05',
             [('YugabyteDB 首輪除錯', {'bold': True, 'color': ORANGE})],
             '處理 BenchmarkSQL、bulk load、RF / schema packing 與 HAProxy 問題',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-05-06～05-14',
             [('三資料庫對標成形', {'bold': True, 'color': ORANGE})],
             '納入 TiDB、CockroachDB、YugabyteDB，統一結果結構與 go-tpc 工具鏈',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
          ],
          col_widths=[2.2, 2.2, 6.5, 1.2])
add_subtitle(s, '此階段完成的是「可比較的工程框架」，早期測試數字不直接納入 v4.7 正式結果。',
             Inches(6.6), color=MUTED, size=14)

# =========================================================================
# Slide 4 — 專案歷程 ②
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 4, 17)
add_title(s, '專案歷程 ② 基準、三節點與治理')
add_subtitle(s, '將測試從可執行提升為可重現、可追溯、可拆解成本來源',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.4),
          header=['時間', '階段', '重大節點', '狀態'],
          rows=[
            ['2026-05-18～05-21',
             [('v4.7 baseline 重構', {'bold': True, 'color': ORANGE})],
             '建立 PoC-DESIGN SSOT、detached suite、gate、marker、summary 與單節點三隔離級對標',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-05-22～06-02',
             [('三節點 controlled experiment', {'bold': True, 'color': ORANGE})],
             '完成 shard × replica × HAProxy 拓樸、12-cell dry-run 與三家 5-cell 結果',
             [('✅ N=1 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-05-20～06-04',
             [('文件與數據治理', {'bold': True, 'color': ORANGE})],
             '建立模板、AI 協作規範、artifact-first 審計與三家 pipeline-log 對齊',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-06-06～06-07',
             [('Phase isolation', {'bold': True, 'color': ORANGE})],
             '分離 S-BASE、S-K8S、T-THRD、X-CROSS，建立配置宣告、基礎指標確認及指標配置',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
          ],
          col_widths=[2.2, 2.4, 6.3, 1.2])
add_subtitle(s, '正式數字必須能追回測試條件、時間戳、結果檔案與完成標記；缺少來源資料 (測試階段) 時不予納入參考。',
             Inches(6.6), color=MUTED, size=14)

# =========================================================================
# Slide 5 — 專案歷程 ③
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 5, 17)
add_title(s, '專案歷程 ③ Kubernetes 到跨區驗證')
add_subtitle(s, '已完成 Kubernetes 對照與跨區技術路徑；正式跨區效能仍受實測結果約束',
             Inches(1.4), color=MUTED)
yellow = RGBColor(0xF1, 0xC4, 0x0F)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.4),
          header=['時間', '階段', '重大節點', '狀態'],
          rows=[
            ['2026-06-08～06-14',
             [('Kubernetes v4.7', {'bold': True, 'color': ORANGE})],
             '由單 cell dry-run 擴充至三資料庫 × limit/unlimit 六組正式 suite',
             [('✅ 6/6 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-06-08～06-17',
             [('跨區設計與前置開發', {'bold': True, 'color': ORANGE})],
             '建立 GCP VM、六節點部署、placement、WAN、chaos、failover 與 pre-flight 規格',
             [('● 框架完成\n部分僅 dry-run', {'color': yellow, 'bold': True})]],
            ['2026-06-18～06-19',
             [('IDC↔GCP 實際驗證', {'bold': True, 'color': ORANGE})],
             '修正 IaC、gate、防火牆與 YugabyteDB placement，三家完成真六節點前期測試',
             [('✅ smoke 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-06-21～06-22',
             [('Determinism 收斂', {'bold': True, 'color': ORANGE})],
             'W=4 重跑浮動過大，改採同 cluster、freeze/unfreeze、變異量/變異參數 與 W=128 baseline',
             [('● 進行中', {'color': yellow, 'bold': True})]],
          ],
          col_widths=[2.2, 2.4, 6.3, 1.2])
add_subtitle(s, '決策界線：可確認六節點跨區交易路徑可行；正式跨家排序須等 W=128、R2～R5 中位數 / 變異係數與回復流程驗收結論。',
             Inches(6.6), color=MUTED, size=14)

# =========================================================================
# Slide 6 — 三家資料庫導入定位
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 6, 17)
add_title(s, '三家資料庫導入定位')

add_table(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.0),
          header=['資料庫', '第一階段觀察', '導入定位'],
          rows=[
            [[('TiDB', {'bold': True, 'color': ORANGE})],
             'VM 與 K8s 吞吐表現較佳；K8s retention 較高',
             '方便優先進入應用情境對接'],
            [[('CockroachDB', {'bold': True, 'color': ORANGE})],
             '一致性需求較強（SSI 預設）；SI / SSI 模式下 retry 與效能成本明顯',
             [('保留觀察', {'bold': True, 'color': RED_EMPH}),
              ('：應用層需評估 Error Handling 容忍度與交易模式', {})]],
            [[('YugabyteDB', {'bold': True, 'color': ORANGE})],
             'VM + HAProxy 表現為佳；K8s 結果目前不宜直接作為導入結論',
             [('保留觀察', {'bold': True, 'color': RED_EMPH}),
              ('：K8s 部署仍需調校與驗證，VM 路徑可進入評估', {})]],
          ],
          col_widths=[2.0, 5.4, 4.7])

# Footer block
add_text(s, Inches(0.6), Inches(3.85), Inches(12.5), Inches(0.4),
         '同 IDC / 同硬體 基準 — vm-3node-haproxy-3s3r-rc, t=128 mean tpmC',
         size=18, bold=True, color=ORANGE)
add_runs(s, Inches(0.6), Inches(4.3), Inches(12.5), Inches(0.4),
         runs=[('三家架構與元件介紹見後續 Slide 7-9；',
                {'color': MUTED, 'size': 14}),
               ('各節點較重損耗待 W=128 baseline 補',
                {'color': RED_EMPH, 'size': 14})])
add_table(s, Inches(0.6), Inches(4.85), Inches(7), Inches(1.2),
          header=['TiDB', 'YugabyteDB', 'CockroachDB'],
          rows=[
            [[('≈ 26,900', {'bold': True, 'size': 20})],
             [('≈ 15,600', {'bold': True, 'size': 20})],
             [('≈ 15,000', {'bold': True, 'size': 20})]],
          ],
          col_widths=[1, 1, 1], header_size=14, body_size=18)
add_text(s, Inches(0.6), Inches(6.2), Inches(12.5), Inches(0.7),
         '本組數字為 controlled experiment 基準（VM 現況資源 × 3 node 拓樸）；不直接代表各家「拿出來就跑」的生產表現；不具任何採購驗收指標用。\n商業實體 / 授權 / 採購層面議題需另案審查。',
         size=12, color=MUTED)

# =========================================================================
# Slide 7 — YugabyteDB 架構
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 7, 17)
add_title(s, 'YugabyteDB 架構')
s.shapes.add_picture(
    '1_MeetingMinutes/assets/arch-ybdb-1.jpg',
    Inches(1.6), Inches(1.55), width=Inches(10.1), height=Inches(4.0))
add_bullets(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.5),
            [[('YSQL：', {'bold': True, 'color': ORANGE}),
              ('PostgreSQL-compatible SQL API；處理 query planning 與 transaction request', {})],
             [('DocDB：', {'bold': True, 'color': ORANGE}),
              ('distributed document store；負責資料儲存、MVCC、tablet 管理與 Raft 複寫', {})],
             [('YB-TServer：', {'bold': True, 'color': ORANGE}),
              ('資料服務節點；承載 YSQL request path 與 DocDB tablet（SQL 與儲存同在 tserver 內，與 TiDB SQL/storage 分離不同）', {})],
             [('YB-Master：', {'bold': True, 'color': ORANGE}),
              ('叢集 metadata 管理；負責 tablet placement、schema metadata 與 cluster coordination', {})],
             [('架構重點：', {'bold': True, 'color': RED_EMPH}),
              ('tserver 一體式；加節點時 SQL 接收能力與 tablet capacity 一起增加', {})]],
            size=12, color=SLATE)

# =========================================================================
# Slide 8 — CockroachDB 架構
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 8, 17)
add_title(s, 'CockroachDB 架構')
s.shapes.add_picture(
    '1_MeetingMinutes/assets/arch-crdb-1.png',
    Inches(3.5), Inches(1.55), width=Inches(6.3), height=Inches(4.0))
add_bullets(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.5),
            [[('SQL layer：', {'bold': True, 'color': ORANGE}),
              ('每個節點都能接 SQL request、產生 query plan、處理 transaction coordination', {})],
             [('Transactional KV / Distribution：', {'bold': True, 'color': ORANGE}),
              ('SQL 轉成 KV operation；distribution layer 路由到 range 所在節點', {})],
             [('Replication / Raft：', {'bold': True, 'color': ORANGE}),
              ('資料以 range 為單位複寫；replica 數直接影響寫入 quorum 與 commit latency', {})],
             [('Storage：', {'bold': True, 'color': ORANGE}),
              ('每節點同時保存資料並處理查詢；無獨立 SQL node / storage node 分層', {})],
             [('架構重點：', {'bold': True, 'color': RED_EMPH}),
              ('對稱式架構；任一節點同時具備 SQL / txn / distribution / replication / storage', {})]],
            size=12, color=SLATE)

# =========================================================================
# Slide 9 — TiDB 架構
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 9, 17)
add_title(s, 'TiDB 架構')
s.shapes.add_picture(
    '1_MeetingMinutes/assets/arch-tidb-1.png',
    Inches(2.5), Inches(1.55), width=Inches(8.3), height=Inches(4.0))
add_bullets(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.5),
            [[('TiDB server：', {'bold': True, 'color': ORANGE}),
              ('SQL 接收層；解析 SQL、產生執行計畫、處理 transaction coordination（不保存主要資料）', {})],
             [('TiKV：', {'bold': True, 'color': ORANGE}),
              ('分散式 row store；資料以 Region 切分並透過 Raft 複寫', {})],
             [('PD (Placement Driver)：', {'bold': True, 'color': ORANGE}),
              ('叢集 metadata 與排程控制中心；負責 timestamp oracle、Region placement 與負載調度', {})],
             [('TiFlash：', {'bold': True, 'color': ORANGE}),
              ('columnar replica，HTAP / analytical query；本輪 OLTP 測試不以 TiFlash 為主路徑', {})],
             [('架構重點：', {'bold': True, 'color': RED_EMPH}),
              ('SQL 層與儲存層分離；TiDB server 擴 SQL 接收、TiKV 擴儲存與 Raft 複寫', {})]],
            size=12, color=SLATE)

# =========================================================================
# Slide 10 — Commit 路徑 × 跨區延遲稅 × 一致性
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 10, 17)
add_title(s, 'Commit 路徑 × 跨區延遲稅 × 一致性')
add_subtitle(s,
             '分散式 DB 跨區後，延遲稅集中在 commit 路徑上「必跨 WAN 的節點」且設計機制不同',
             Inches(1.4), color=MUTED, size=14)

GREEN_S = RGBColor(0x27, 0xAE, 0x60)
COL_W = Inches(4.0)
COL_Y = Inches(2.05)
COL_H = Inches(2.5)
gap = Inches(0.15)
cols = [
    ('TiDB',
     ('時序', '中央 PD TSO\n每 txn 取 ts +1 RTT'),
     ('Commit', '2PC (Percolator)\nprewrite → commit'),
     ('延遲稅集中', 'PD TSO 跨區')),
    ('CockroachDB',
     ('時序', 'HLC（本地時鐘 + counter）\n無 TSO RTT'),
     ('Commit', 'gateway 為 txn coord\nintent → resolve'),
     ('延遲稅集中', 'leaseholder + intent resolve')),
    ('YugabyteDB',
     ('時序', 'HLC（類 CRDB）\n無 TSO RTT'),
     ('Commit', 'YSQL → DocDB tablet\nraft + 狀態表'),
     ('延遲稅集中', 'YSQL ↔ DocDB IPC\n+ tablet leader 跨區')),
]
for i, (name, *rows) in enumerate(cols):
    x = Inches(0.6) + (COL_W + gap) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, COL_Y, COL_W, COL_H)
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x26, 0x2F, 0x46)
    card.line.color.rgb = ACCENT_RED; card.line.width = Pt(1)
    card.shadow.inherit = False
    add_text(s, x + Inches(0.18), COL_Y + Inches(0.08), COL_W - Inches(0.36), Inches(0.35),
             name, size=18, bold=True, color=ORANGE)
    yp = COL_Y + Inches(0.55)
    for label, val in rows:
        add_text(s, x + Inches(0.18), yp, COL_W - Inches(0.36), Inches(0.22),
                 label + '：', size=11, bold=True, color=ACCENT_RED)
        add_text(s, x + Inches(0.18), yp + Inches(0.22), COL_W - Inches(0.36), Inches(0.45),
                 val, size=10, color=SLATE)
        yp += Inches(0.62)

add_text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4),
         '降稅的兩個關鍵動作', size=16, bold=True, color=ORANGE)
add_bullets(s, Inches(0.8), Inches(5.1), Inches(12.3), Inches(1.4),
            [[('Placement P-A：', {'bold': True, 'color': ORANGE}),
              ('把 voter 多數拉回 IDC，quorum 不跨 WAN → 寫入稅從 RTT 降到 0', {})],
             [('Follower / stale read：', {'bold': True, 'color': ORANGE}),
              ('讀本地 follower，bounded staleness 換 0 RTT → 適用 A-A-RO workload', {})],
             [('YugabyteDB 多一筆稅：', {'bold': True, 'color': RED_EMPH}),
              ('YSQL postgres backend ↔ DocDB tablet 是兩 process；跨區 IPC + RTT 雙重放大', {})]],
            size=13, color=SLATE)
add_runs(s, Inches(0.6), Inches(6.65), Inches(12.5), Inches(0.4),
         runs=[('⚠ 強一致與跨區效能無法同時最佳化',
                {'bold': True, 'color': RGBColor(0xF1, 0xC4, 0x0F), 'size': 13}),
               (' — 必須選 workload 模式（A-S / A-A-RO / A-A）',
                {'color': MUTED, 'size': 13})])

# =========================================================================
# Slide 11 — Range vs Hash Sharding
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 11, 17)
add_title(s, 'Range vs Hash Sharding')
add_subtitle(s,
             '「row 屬哪個 shard」的兩種演算法 — 直接影響 locality / hotspot / range scan 效能',
             Inches(1.4), color=MUTED, size=14)

YELLOW_SH = RGBColor(0xF1, 0xC4, 0x0F)
card_w = Inches(6.0); card_h = Inches(2.6); card_y = Inches(2.05)
for i, (title, color, items) in enumerate([
    ('Range Sharding', GREEN_S,
     [('演算法', '依 key 區間切（id 1-1000 → shard-1）'),
      ('優點', '範圍掃描快、locality 高、index/排序自然對齊'),
      ('缺點', 'Sequential write (auto_inc / timestamp) 集中熱點'),
      ('代表 DB', 'TiDB / CockroachDB / YBDB secondary index')]),
    ('Hash Sharding', ORANGE,
     [('演算法', 'hash(key) % N → shard'),
      ('優點', '寫入分散、無熱點、shard 容量均勻'),
      ('缺點', '範圍掃描需 fan-out 全 shard；跨 shard JOIN 貴'),
      ('代表 DB', 'YugabyteDB (YSQL PK 預設) / Cassandra')]),
]):
    x = Inches(0.6) + (card_w + Inches(0.2)) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x26, 0x2F, 0x46)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    card.shadow.inherit = False
    add_text(s, x + Inches(0.2), card_y + Inches(0.1), card_w - Inches(0.4), Inches(0.4),
             title, size=18, bold=True, color=color)
    yp = card_y + Inches(0.55)
    for label, val in items:
        add_runs(s, x + Inches(0.25), yp, card_w - Inches(0.5), Inches(0.4),
                 runs=[(label + '：', {'bold': True, 'color': color, 'size': 12}),
                       (val, {'color': SLATE, 'size': 12})])
        yp += Inches(0.48)

add_table(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.6),
          header=['觀察', '解釋'],
          rows=[
            ['TiDB tpmC ≈ 26,900（最高）',
             'Range + warehouse-aware split + leader balance 對 OLTP locality 最友善'],
            ['CockroachDB tpmC ≈ 15,000',
             'Range 但 SSI 預設 + retry 成本拉低吞吐'],
            ['YugabyteDB tpmC ≈ 15,600 (VM)',
             'Hash 預設打散；無熱點但 cross-tablet txn 比例高，吞吐受限於 commit fanout'],
          ],
          col_widths=[4, 8.1], header_size=13, body_size=12)

add_runs(s, Inches(0.6), Inches(6.75), Inches(12.5), Inches(0.35),
         runs=[('TPC-C warehouse-centric workload',
                {'bold': True, 'color': ORANGE, 'size': 12}),
               (' 偏好 range（同 warehouse row co-located）；'
                'hash 在 hot warehouse 場景反能避開單 shard 飽和',
                {'color': SLATE, 'size': 12})])

# =========================================================================
# Slide 12 — VM 與 Kubernetes 效能差異
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 12, 17)
add_title(s, 'VM 與 Kubernetes 效能差異')
add_subtitle(s, 'K8s 化對 TiDB 與 CockroachDB 屬可接受範圍；YugabyteDB 在 K8s 下退化顯著，列為調校與驗證項',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.5),
          header=['資料庫', 'S-BASE (VM)\ntpmC', 'S-K8S unlimit\ntpmC', 'S-K8S limit\ntpmC',
                  'unlimit / VM', 'limit / VM', 'p99 unlimit Δ', 'p99 limit Δ'],
          rows=[
            [[('TiDB v8.5.2', {'bold': True})],
             '26,947', '23,442.9', '15,751.9',
             [('87.0 %', {'bold': True, 'color': RED_EMPH})],
             '58.5 %', '+17 %', '+111 %'],
            [[('CockroachDB v26.2', {'bold': True})],
             '15,033', '12,196.7', '6,493.5',
             [('81.1 %', {'bold': True, 'color': RED_EMPH})],
             '43.2 %', '+27 %', '+192 %'],
            [[('YugabyteDB 2025.2', {'bold': True})],
             '15,632', '2,997.6', '1,604.5',
             [('19.2 % ⚠', {'bold': True, 'color': RED_EMPH})],
             '10.3 % ⚠', '+669 %', '+1556 %'],
          ],
          col_widths=[2.2, 1.5, 1.5, 1.5, 1.4, 1.2, 1.4, 1.4],
          header_size=12, body_size=13)
add_text(s, Inches(0.6), Inches(4.85), Inches(12), Inches(0.5),
         '觀察解讀', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(0.8), Inches(5.3), Inches(12), Inches(1.6),
            ['TiDB / CockroachDB：K8s unlimit 保留率 81-87 %，屬可接受範圍；limit 情境大幅衰退反映資源壓縮',
             [('YugabyteDB：K8s unlimit 僅 19.2 %、p99 +669 %。成因推論：YSQL + DocDB 雙 process 在 K8s pod IPC + CPU contention 放大。',
               {}),
              ('成因未定位，列後續調校項', {'bold': True, 'color': RED_EMPH})],
             'S-BASE ↔ S-K8S 可直接比較（同 workload W=128 t=128）'],
            size=15, color=SLATE)

# =========================================================================
# Slide 13 — 跨區域 / 跨專線進度
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 13, 17)
add_title(s, '跨區域 / 跨專線進度')
add_subtitle(s,
             'IaC、playbook、suite scripts、chrony drift gate 已就位；正式量測待 W=128 + Determinism gate',
             Inches(1.4), color=MUTED, size=14)

add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
         '✅ 已完成', size=20, bold=True, color=RGBColor(0x27, 0xAE, 0x60))
add_bullets(s, Inches(0.9), Inches(2.5), Inches(12), Inches(2.2),
            ['技術議題定案（GCP 5 VM 拓樸、placement P-A/P-B、Chaos Engineering 等）',
             'IaC、playbook、suite scripts、chrony drift gate（drift_median 0.017 ms；時間偏移 between IDC & GCP）',
             '06-18〜06-19：修正 IaC / 防火牆 / YugabyteDB placement，三家 DB / 六節點前期測試驗證',
             '06-21〜06-22：Determinism v2 — 同 cluster R1-R5、scheduler / balancer freeze'],
            size=16, color=SLATE)

add_text(s, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
         '● 待執行', size=20, bold=True, color=yellow)
add_bullets(s, Inches(0.9), Inches(5.3), Inches(12), Inches(1.6),
            ['調度 scheduler (開/關)、round-only runner、Warmup、Placement (replica / lease 位置)、測試結果彙整',
             '回到 W=128 正式基準量測',
             'Failover 與 Chaos Engineering 實際驗證'],
            size=16, color=SLATE)

# =========================================================================
# Slide 14 — P-A / P-B Placement 解說
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 14, 17)
add_title(s, 'P-A / P-B Placement — 跨區寫入路徑的兩種設定')
add_subtitle(s,
             '6-node cluster (3 IDC + 3 GCP) 下，voter 與 leader 的分布決定 raft commit 是否走 WAN',
             Inches(1.4), color=MUTED, size=14)

GREEN_PB = RGBColor(0x27, 0xAE, 0x60)
YELLOW_PB = RGBColor(0xF1, 0xC4, 0x0F)
BOX_Y2 = Inches(2.0)
BOX_H2 = Inches(3.4)
MINI_W2 = Inches(2.45)
MINI_H2 = Inches(0.7)

# --- P-A card ---
PA_X2, PA_W2 = Inches(0.6), Inches(6.0)
card_a = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PA_X2, BOX_Y2, PA_W2, BOX_H2)
card_a.fill.solid(); card_a.fill.fore_color.rgb = RGBColor(0x26, 0x2F, 0x46)
card_a.line.color.rgb = GREEN_PB; card_a.line.width = Pt(1.5)
card_a.shadow.inherit = False
add_text(s, PA_X2 + Inches(0.2), BOX_Y2 + Inches(0.1), PA_W2 - Inches(0.4), Inches(0.4),
         'P-A · 多數 voter 在 IDC', size=18, bold=True, color=GREEN_PB)
add_text(s, PA_X2 + Inches(0.2), BOX_Y2 + Inches(0.55), PA_W2 - Inches(0.4), Inches(0.3),
         '情境：正常營運（IDC 為主，GCP 為備）', size=12, color=MUTED)

idc_x2 = PA_X2 + Inches(0.4); diag_y2 = BOX_Y2 + Inches(1.05)
idc_b2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, idc_x2, diag_y2, MINI_W2, MINI_H2)
idc_b2.fill.solid(); idc_b2.fill.fore_color.rgb = RGBColor(0x1A, 0x4A, 0x2F)
idc_b2.line.color.rgb = GREEN_PB; idc_b2.shadow.inherit = False
add_text(s, idc_x2, diag_y2, MINI_W2, MINI_H2,
         'IDC\n2 voter/shard · leader',
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
gcp_x2 = idc_x2 + MINI_W2 + Inches(0.15)
gcp_b2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, gcp_x2, diag_y2, MINI_W2, MINI_H2)
gcp_b2.fill.solid(); gcp_b2.fill.fore_color.rgb = RGBColor(0x33, 0x3D, 0x55)
gcp_b2.line.color.rgb = MUTED; gcp_b2.shadow.inherit = False
add_text(s, gcp_x2, diag_y2, MINI_W2, MINI_H2,
         'GCP\n1 voter/shard · follower',
         size=12, color=SLATE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, PA_X2 + Inches(0.25), BOX_Y2 + Inches(2.0),
            PA_W2 - Inches(0.5), Inches(1.3),
            ['Quorum 在 IDC，GCP 不擋 commit',
             '適用 workload：A-S（IDC main / GCP standby）',
             [('預期 tpmC 衝擊：', {}),
              ('約 −10〜30 %', {'bold': True, 'color': ORANGE})]],
            size=13, color=SLATE, bullet_color=GREEN_PB)

# --- P-B card ---
PB_X2, PB_W2 = Inches(7.0), Inches(5.8)
card_b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PB_X2, BOX_Y2, PB_W2, BOX_H2)
card_b2.fill.solid(); card_b2.fill.fore_color.rgb = RGBColor(0x26, 0x2F, 0x46)
card_b2.line.color.rgb = YELLOW_PB; card_b2.line.width = Pt(1.5)
card_b2.shadow.inherit = False
add_text(s, PB_X2 + Inches(0.2), BOX_Y2 + Inches(0.1), PB_W2 - Inches(0.4), Inches(0.4),
         'P-B · leader 散在兩區', size=18, bold=True, color=YELLOW_PB)
add_text(s, PB_X2 + Inches(0.2), BOX_Y2 + Inches(0.55), PB_W2 - Inches(0.4), Inches(0.3),
         '情境：退化 / failover 後 / chaos 測試', size=12, color=MUTED)

idc3_x = PB_X2 + Inches(0.3)
idc3_b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, idc3_x, diag_y2, Inches(2.35), MINI_H2)
idc3_b.fill.solid(); idc3_b.fill.fore_color.rgb = RGBColor(0x4A, 0x3A, 0x1A)
idc3_b.line.color.rgb = YELLOW_PB; idc3_b.shadow.inherit = False
add_text(s, idc3_x, diag_y2, Inches(2.35), MINI_H2,
         'IDC\nleader + voter + arbiter',
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
gcp3_x = idc3_x + Inches(2.5)
gcp3_b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, gcp3_x, diag_y2, Inches(2.35), MINI_H2)
gcp3_b.fill.solid(); gcp3_b.fill.fore_color.rgb = RGBColor(0x4A, 0x3A, 0x1A)
gcp3_b.line.color.rgb = YELLOW_PB; gcp3_b.shadow.inherit = False
add_text(s, gcp3_x, diag_y2, Inches(2.35), MINI_H2,
         'GCP\nleader + voter + arbiter',
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, idc3_x, diag_y2 + MINI_H2 + Inches(0.02),
         Inches(2.35) * 2 + Inches(0.15), Inches(0.2),
         '↔ WAN 必擋 raft commit', size=10, color=YELLOW_PB,
         align=PP_ALIGN.CENTER)

add_bullets(s, PB_X2 + Inches(0.25), BOX_Y2 + Inches(2.0),
            PB_W2 - Inches(0.5), Inches(1.3),
            ['寫入需 WAN quorum，commit latency = RTT',
             '適用 workload：A-A（雙寫）/ A-A-RO（IDC write, GCP read）',
             [('預期 tpmC 衝擊：', {}),
              ('約 −30〜60 %', {'bold': True, 'color': RED_EMPH})]],
            size=13, color=SLATE, bullet_color=YELLOW_PB)

# Bottom takeaway
add_runs(s, Inches(0.6), Inches(5.7), Inches(12.5), Inches(0.4),
         runs=[('對標重點：', {'bold': True, 'color': ORANGE, 'size': 14}),
               ('P-A 量「正常營運下的跨區 cost」；P-B 量「退化/failover 後的跨區 cost」',
                {'color': SLATE, 'size': 14})])
add_runs(s, Inches(0.6), Inches(6.15), Inches(12.5), Inches(0.4),
         runs=[('兩者相減 ≈ failover 觸發時的額外吞吐損失',
                {'color': MUTED, 'size': 13})])

# =========================================================================
# Slide 15 — X-CROSS 初步結果
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 15, 17)
add_title(s, 'X-CROSS 初步結果 (06-22〜06-23)')
add_subtitle(s, 'W=4 framework 驗證；不與 S-BASE/S-K8S 直接對比 · 同 cluster 5 rounds、W=4、16 threads、每 round 5 min、controller = .31 (IDC client)',
             Inches(1.4), color=MUTED, size=13)
green_check = RGBColor(0x27, 0xAE, 0x60)
add_table(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(2.1),
          header=['資料庫', 'R1', 'R2', 'R3', 'R4', 'R5', '中位數 (有效輪)', 'CV'],
          rows=[
            [[('TiDB v8.5.2', {'bold': True})],
             '9,525.5', '9,553.2', '9,786.9', '9,393.2', '9,530.8',
             [('9,566.0 (R2-R5)', {'bold': True, 'color': RED_EMPH})],
             [('1.67 % ✅', {'bold': True, 'color': green_check})]],
            [[('CockroachDB v26.2', {'bold': True})],
             '8,409.5', '8,055.3', '7,902.5', '7,720.9', '7,472.3',
             [('7,787.8 (R2-R5)', {'bold': True, 'color': RED_EMPH})],
             [('3.23 % ✅', {'bold': True, 'color': green_check})]],
            [[('YugabyteDB 2025.2', {'bold': True})],
             '102.0', '226.9', '6,424.2', '6,259.3', '6,206.2',
             [('6,296.6 (R3-R5)', {'bold': True, 'color': RED_EMPH})],
             [('1.82 % ✅', {'bold': True, 'color': green_check})]],
          ],
          col_widths=[2.4, 1.1, 1.1, 1.1, 1.1, 1.1, 2.3, 1.3],
          header_size=12, body_size=12)
add_text(s, Inches(0.6), Inches(4.35), Inches(12), Inches(0.4),
         '已確認結論', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(0.8), Inches(4.8), Inches(12.3), Inches(1.7),
            ['06-21 觀察到的 ±526 % / ±50 % 變異主因為「每輪重新部署」造成 placement / cache / scheduler 狀態異動，與 W=4 contention 為兩個獨立來源',
             '同一 cluster 連跑時，三家 W=4 R2-R5（或 YBDB R3-R5）的 CV ≤ 5 %，重現性已建立',
             'YBDB Idle=0 解法：gate 改用 get_load_move_completion=100%；timed run 前 set_load_balancer_enabled=0',
             'CRDB lease gate SQL 改用 SHOW RANGES FROM DATABASE tpcc WITH TABLES, DETAILS（v26.2 相容）'],
            size=13, color=SLATE)
add_runs(s, Inches(0.6), Inches(6.6), Inches(12.5), Inches(0.4),
         runs=[('⚠ Caveat：', {'bold': True, 'color': yellow, 'size': 13}),
               ('W=4 framework / contention 驗證 ≠ 正式基準；R2-R5（YBDB R3-R5）CV ≤ 5 % 已建立但',
                {'color': MUTED, 'size': 13}),
               ('不可作跨家排名', {'bold': True, 'color': RED_EMPH, 'size': 13}),
               ('；後續需回 W=128。', {'color': MUTED, 'size': 13})])

# =========================================================================
# Slide 16 — 建議決策框架
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 16, 17)
add_title(s, '建議決策框架')
add_subtitle(s, '不是選一家，是先界定 application 條件、再排序候選；僅提供做決策框架參考依據用',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.0),
          header=['分類', '內容', '處置'],
          rows=[
            [[('短期候選', {'bold': True, 'color': RGBColor(0x27, 0xAE, 0x60)})],
             [('TiDB', {'bold': True})],
             '優先進入應用情境對接，無門檻'],
            [[('保留觀察', {'bold': True, 'color': ORANGE})],
             'CockroachDB',
             '依一致性需求、Error Handling 容忍度、維運成本評估'],
            [[('保留觀察', {'bold': True, 'color': ORANGE})],
             'YugabyteDB',
             'VM 路徑可評估；K8s 路徑需先完成部署層級調校與驗證'],
            [[('暫不作結論', {'bold': True, 'color': RED_EMPH})],
             '跨區域場景、K8s 退化未定位項',
             '待下一階段 W=128 測試數據與調校產出再回頭評估'],
          ],
          col_widths=[1.8, 3.5, 6.8])
add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
         'Application 需要確認的議題', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(0.8), Inches(4.8), Inches(12.3), Inches(2.3),
            [[('1. 交易一致性需求：', {'bold': True, 'color': ORANGE}),
              ('是否需要 SERIALIZABLE / SSI？是否可接受 READ COMMITTED & 產品設計架構轉 CAP 的複雜度及可行性？', {})],
             [('2. 可接受延遲：', {'bold': True, 'color': ORANGE}),
              ('在尖峰 t=128 等級負載下，p99 是否需要 500ms / 1s / 2s 以內？', {})],
             [('3. retry / timeout 行為：', {'bold': True, 'color': ORANGE}),
              ('應用層是否能接受 SI / SSI 模式下的 retry 機制？', {})],
             [('4. RTO / RPO：', {'bold': True, 'color': ORANGE}),
              ('跨區域 failover 是否需要 < 30s RTO？是否容許資料 lag？', {})],
             [('5. 連線層 / 交易模式調整：', {'bold': True, 'color': ORANGE}),
              ('是否能配合 HAProxy / PgBouncer / 連線池與短交易模式？', {})]],
            size=13, color=SLATE)

# =========================================================================
# Slide 17 — 後續推進階段
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 17, 17)
add_title(s, '後續推進階段')
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.1),
          header=['短期 EX: Y26/07', '中期', '決策'],
          rows=[
            ['Cross Region W=4 框架已完成初步流程驗收\n啟動正式 W=128 基準量測',
             'A-S / A-A-RO / A-A 三個 workload profile × P-A/P-B placement × 三家資料庫，依檢驗指標分批 / 量測',
             'Application owner 完成五項議題確認'],
            ['W=4 框架已驗，接 W=128 基準',
             'Failover 與 Chaos Engineering 測試',
             '原廠後勤對接狀況說明'],
            ['freeze/unfreeze、round-only runner、暖機、placement gate 全部通過 審查',
             '跨區域 analytics 第二份報告',
             ''],
            ['X-CROSS 完成 determinism 流程 驗收，啟動正式 W=128 基準 量測',
             '補齊 YugabyteDB K8s 退化成因調查（P-A/P-B、A-S/A-A-RO/A-A）',
             ''],
          ],
          col_widths=[4.5, 5.0, 2.6],
          header_size=14, body_size=13)

# =========================================================================
# Save
# =========================================================================
out = '/Users/wn.lin/vscode-git/dba_career/poc/1_MeetingMinutes/0626-slide-v6.pptx'
prs.save(out)
print(f'wrote {out}')
