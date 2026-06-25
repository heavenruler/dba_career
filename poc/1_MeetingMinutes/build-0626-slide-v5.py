#!/usr/bin/env python3
"""Build 0626-slide-v5.pptx — editable PowerPoint matching Marp dark-navy theme.

Produces text-box / table based slides (not flattened images) so user can edit.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---- palette ------------------------------------------------------------
BG_NAVY     = RGBColor(0x1E, 0x25, 0x38)
ACCENT_RED  = RGBColor(0xE7, 0x4C, 0x3C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLATE       = RGBColor(0xE8, 0xEA, 0xED)
MUTED       = RGBColor(0xA0, 0xA8, 0xB8)
ORANGE      = RGBColor(0xF3, 0x9C, 0x12)
RED_EMPH    = RGBColor(0xB0, 0x3A, 0x14)
TABLE_BG    = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ALT   = RGBColor(0xF1, 0xF3, 0xF6)
TABLE_HEAD  = RGBColor(0xD8, 0xDD, 0xE6)
TABLE_TEXT  = RGBColor(0x11, 0x14, 0x18)
TABLE_BORD  = RGBColor(0x8C, 0x95, 0xA3)

FONT_CJK    = 'PingFang TC'
FONT_LATIN  = 'Inter'

# ---- layout -------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, *, size=22, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font or FONT_CJK
    return box

def add_runs(slide, x, y, w, h, runs, *, size=20, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """runs = list of (text, {bold,color,size,font}) -- all on one paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, opts in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.color.rgb = opts.get('color', SLATE)
        r.font.name = opts.get('font', FONT_CJK)
    return box

def add_bullets(slide, x, y, w, h, items, *, size=18, color=SLATE,
                bullet_color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # XML-level bullet char
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgb = etree.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', f'{(bullet_color or ACCENT_RED).rgb:06X}' if hasattr((bullet_color or ACCENT_RED), 'rgb') else 'E74C3C')
        pPr.set('indent', '-228600')
        pPr.set('marL', '228600')

        if isinstance(item, str):
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = FONT_CJK
        else:
            for txt, opts in item:
                r = p.add_run()
                r.text = txt
                r.font.size = Pt(opts.get('size', size))
                r.font.bold = opts.get('bold', False)
                r.font.color.rgb = opts.get('color', color)
                r.font.name = opts.get('font', FONT_CJK)
    return box

def add_header_footer(slide, page_num, total):
    # top accent bar
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Pt(6), ACCENT_RED)
    # header text
    add_text(slide, Inches(0.4), Inches(0.18), Inches(6), Inches(0.3),
             'PoC 第一階段成果與下一步決策', size=12, color=MUTED)
    # footer
    add_text(slide, Inches(0.4), Inches(7.05), Inches(6), Inches(0.3),
             '2026-06-26 · DBA', size=12, color=MUTED)
    # page number
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3),
             str(page_num), size=12, color=MUTED, align=PP_ALIGN.RIGHT)

def add_title(slide, text, y=Inches(0.55)):
    add_text(slide, Inches(0.6), y, Inches(12), Inches(0.7),
             text, size=30, bold=True, color=WHITE)
    # underline
    add_rect(slide, Inches(0.6), y + Inches(0.7), Inches(12), Pt(2), ACCENT_RED)

def add_subtitle(slide, text, y, *, color=MUTED, size=18):
    """Subtitle with red left bar."""
    bar_x = Inches(0.6)
    bar_w = Pt(3)
    add_rect(slide, bar_x, y, bar_w, Inches(0.4), ACCENT_RED)
    add_text(slide, bar_x + Emu(80000), y, Inches(12), Inches(0.4),
             text, size=size, color=color, anchor=MSO_ANCHOR.MIDDLE)

# ---- table builder ------------------------------------------------------
def style_cell(cell, *, bg=TABLE_BG, color=TABLE_TEXT, bold=False, size=14,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, runs=None):
    """Style a table cell. runs=[(text,opts)...] overrides default text rendering."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.margin_left = cell.margin_right = Inches(0.08)
    cell.margin_top = cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = anchor
    tf = cell.text_frame
    tf.word_wrap = True
    # clear existing text
    for p in list(tf.paragraphs):
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
    p = tf.paragraphs[0]
    p.alignment = align
    if runs is None:
        return  # caller will set text manually via cell.text
    for txt, opts in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', bold)
        r.font.color.rgb = opts.get('color', color)
        r.font.name = opts.get('font', FONT_CJK)

def add_table(slide, x, y, w, h, header, rows, *, col_widths=None,
              header_size=14, body_size=13):
    """header = [str]; rows = [[cell_value]] where cell_value is str or [(text,opts)..]."""
    n_rows = len(rows) + 1
    n_cols = len(header)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    # column widths
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    # header row
    for i, txt in enumerate(header):
        cell = tbl.cell(0, i)
        style_cell(cell, bg=TABLE_HEAD, color=TABLE_TEXT, bold=True,
                   size=header_size, align=PP_ALIGN.CENTER,
                   runs=[(txt, {'bold': True, 'size': header_size})])
    # body rows
    for r_idx, row in enumerate(rows):
        bg = TABLE_BG if r_idx % 2 == 0 else TABLE_ALT
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            if isinstance(val, str):
                style_cell(cell, bg=bg, size=body_size,
                           runs=[(val, {'size': body_size})])
            else:
                # list of (text, opts)
                style_cell(cell, bg=bg, size=body_size, runs=val)
    return tbl

# =========================================================================
# Slide 1 — Title
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_rect(s, Emu(0), Emu(0), SLIDE_W, Pt(6), ACCENT_RED)
add_text(s, Inches(0.4), Inches(0.18), Inches(6), Inches(0.3),
         'PoC 第一階段成果與下一步決策', size=12, color=MUTED)
add_text(s, Inches(0), Inches(3.0), SLIDE_W, Inches(0.9),
         '第一階段 PoC 成果與下一步決策', size=44, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.2), SLIDE_W, Inches(0.5),
         '分散式資料庫架構 × 跨區域驗證 × 決策框架', size=20, color=MUTED,
         align=PP_ALIGN.CENTER)
# divider
add_rect(s, Inches(5), Inches(5.0), Inches(3.3), Pt(2), ACCENT_RED)
add_text(s, Inches(0), Inches(5.2), SLIDE_W, Inches(0.5),
         '2026-06-26    |    DBA', size=18, color=SLATE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3),
         '1', size=12, color=MUTED, align=PP_ALIGN.RIGHT)

# =========================================================================
# Slide 2 — Outline
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 2, 11)
add_title(s, 'Outline')
add_bullets(s, Inches(0.8), Inches(1.7), Inches(10), Inches(5),
            ['專案迄今歷程',
             '第一階段 PoC 測試數據彙整',
             '第二階段 跨區域 / 跨專線執行進度',
             '決策框架說明',
             '後續推進'],
            size=22, color=SLATE, bullet_color=ACCENT_RED)

# =========================================================================
# Slide 3 — 專案歷程 ①
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 3, 11)
add_title(s, '專案歷程 ① 研究定義到跨家框架')
add_subtitle(s, '先定義問題與比較口徑，再建立可重複執行的基礎設施與三資料庫共同工具鏈',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.4),
          header=['時間', '階段', '重大節點', '狀態'],
          rows=[
            ['2026-03-30～04-10',
             [('前期研究', {'bold': True, 'color': ORANGE})],
             '定義分散式 SQL、跨區同鍵寫入、follower read、HA/DR 與九項 survey 評估面向',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-04-21～04-27',
             [('IaC 與第一版測試鏈', {'bold': True, 'color': ORANGE})],
             [('建立多測項部署、HAProxy、VM / Kubernetes 流程及獨立壓測 client\n', {}),
              ('{FIXME}', {'bold': True, 'color': RED_EMPH}),
              (' Chain 的 flow & Key Point', {})],
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-04-28～05-05',
             [('YugabyteDB 首輪除錯', {'bold': True, 'color': ORANGE})],
             '處理 BenchmarkSQL、bulk load、RF / schema packing 與 HAProxy 問題',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-05-06～05-14',
             [('三資料庫對標成形', {'bold': True, 'color': ORANGE})],
             '納入 TiDB、CockroachDB、YugabyteDB，統一結果結構與 go-tpc 工具鏈',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
          ],
          col_widths=[2.2, 2.2, 6.5, 1.2])
add_subtitle(s, '此階段完成的是「可比較的工程框架」，早期測試數字不直接納入 v4.7 正式結果。',
             Inches(6.6), color=MUTED, size=14)

# =========================================================================
# Slide 4 — 專案歷程 ②
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 4, 11)
add_title(s, '專案歷程 ② 基準、三節點與治理')
add_subtitle(s, '將測試從可執行提升為可重現、可追溯、可拆解成本來源',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.4),
          header=['時間', '階段', '重大節點', '狀態'],
          rows=[
            ['2026-05-18～05-21',
             [('v4.7 baseline 重構', {'bold': True, 'color': ORANGE})],
             '建立 PoC-DESIGN SSOT、detached suite、gate、marker、summary 與單節點三隔離級對標',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-05-22～06-02',
             [('三節點 controlled experiment', {'bold': True, 'color': ORANGE})],
             '完成 shard × replica × HAProxy 拓樸、12-cell dry-run 與三家 5-cell 結果',
             [('✅ N=1 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-05-20～06-04',
             [('文件與數據治理', {'bold': True, 'color': ORANGE})],
             '建立模板、AI 協作規範、artifact-first 審計與三家 pipeline-log 對齊',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-06-06～06-07',
             [('Phase isolation', {'bold': True, 'color': ORANGE})],
             '分離 S-BASE、S-K8S、T-THRD、X-CROSS，建立配置宣告、基礎指標確認及指標配置',
             [('✅ 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
          ],
          col_widths=[2.2, 2.4, 6.3, 1.2])
add_subtitle(s, '正式數字必須能追回測試條件、時間戳、結果檔案與完成標記；缺少來源資料 (測試階段) 時不予納入參考。',
             Inches(6.6), color=MUTED, size=14)

# =========================================================================
# Slide 5 — 專案歷程 ③
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 5, 11)
add_title(s, '專案歷程 ③ Kubernetes 到跨區驗證')
add_subtitle(s, '已完成 Kubernetes 對照與跨區技術路徑；正式跨區效能仍受實測結果約束',
             Inches(1.4), color=MUTED)
yellow = RGBColor(0xF1, 0xC4, 0x0F)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.4),
          header=['時間', '階段', '重大節點', '狀態'],
          rows=[
            ['2026-06-08～06-14',
             [('Kubernetes v4.7', {'bold': True, 'color': ORANGE})],
             '由單 cell dry-run 擴充至三資料庫 × limit/unlimit 六組正式 suite',
             [('✅ 6/6 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-06-08～06-17',
             [('跨區設計與前置開發', {'bold': True, 'color': ORANGE})],
             '建立 GCP VM、六節點部署、placement、WAN、chaos、failover 與 pre-flight 規格',
             [('● 框架完成\n部分僅 dry-run', {'color': yellow, 'bold': True})]],
            ['2026-06-18～06-19',
             [('IDC↔GCP 實際驗證', {'bold': True, 'color': ORANGE})],
             '修正 IaC、gate、防火牆與 YugabyteDB placement，三家完成真六節點前期測試',
             [('✅ smoke 完成', {'color': RGBColor(0x27, 0xAE, 0x60), 'bold': True})]],
            ['2026-06-21～06-22',
             [('Determinism 收斂', {'bold': True, 'color': ORANGE})],
             'W=4 重跑浮動過大，改採同 cluster、freeze/unfreeze、變異量/變異參數 與 W=128 baseline',
             [('● 進行中', {'color': yellow, 'bold': True})]],
          ],
          col_widths=[2.2, 2.4, 6.3, 1.2])
add_subtitle(s, '決策界線：可確認六節點跨區交易路徑可行；正式跨家排序須等 W=128、R2～R5 中位數 / 變異係數與回復流程驗收結論。',
             Inches(6.6), color=MUTED, size=14)

# =========================================================================
# Slide 6 — 三家資料庫導入定位
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 6, 11)
add_title(s, '三家資料庫導入定位')

add_table(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.0),
          header=['資料庫', '第一階段觀察', '導入定位'],
          rows=[
            [[('TiDB', {'bold': True, 'color': ORANGE})],
             'VM 與 K8s 吞吐表現較佳；K8s retention 較高',
             '方便優先進入應用情境對接'],
            [[('CockroachDB', {'bold': True, 'color': ORANGE})],
             '一致性需求較強（SSI 預設）；SI / SSI 模式下 retry 與效能成本明顯',
             [('保留觀察', {'bold': True, 'color': RED_EMPH}),
              ('：應用層需評估 Error Handling 容忍度與交易模式', {})]],
            [[('YugabyteDB', {'bold': True, 'color': ORANGE})],
             'VM + HAProxy 表現為佳；K8s 結果目前不宜直接作為導入結論',
             [('保留觀察', {'bold': True, 'color': RED_EMPH}),
              ('：K8s 部署仍需調校與驗證，VM 路徑可進入評估', {})]],
          ],
          col_widths=[2.0, 5.4, 4.7])

# Footer block
add_text(s, Inches(0.6), Inches(3.85), Inches(12.5), Inches(0.4),
         '同 IDC / 同硬體 基準 — vm-3node-haproxy-3s3r-rc, t=128 mean tpmC',
         size=18, bold=True, color=ORANGE)
add_runs(s, Inches(0.6), Inches(4.3), Inches(12.5), Inches(0.4),
         runs=[('{FIXME}', {'bold': True, 'color': RED_EMPH, 'size': 14}),
               (' 測試架構圖補完 ; db component intro ; 各節點較重損耗',
                {'color': MUTED, 'size': 14})])
add_table(s, Inches(0.6), Inches(4.85), Inches(7), Inches(1.2),
          header=['TiDB', 'YugabyteDB', 'CockroachDB'],
          rows=[
            [[('≈ 26,900', {'bold': True, 'size': 20})],
             [('≈ 15,600', {'bold': True, 'size': 20})],
             [('≈ 15,000', {'bold': True, 'size': 20})]],
          ],
          col_widths=[1, 1, 1], header_size=14, body_size=18)
add_text(s, Inches(0.6), Inches(6.2), Inches(12.5), Inches(0.7),
         '本組數字為 controlled experiment 基準（VM 現況資源 × 3 node 拓樸）；不直接代表各家「拿出來就跑」的生產表現；不具任何採購驗收指標用。\n商業實體 / 授權 / 採購層面議題需另案審查。',
         size=12, color=MUTED)

# =========================================================================
# Slide 7 — VM 與 Kubernetes 效能差異
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 7, 11)
add_title(s, 'VM 與 Kubernetes 效能差異')
add_subtitle(s, 'K8s 化對 TiDB 與 CockroachDB 屬可接受範圍；YugabyteDB 在 K8s 下退化顯著，列為調校與驗證項',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.5),
          header=['資料庫', 'S-BASE (VM)\ntpmC', 'S-K8S unlimit\ntpmC', 'S-K8S limit\ntpmC',
                  'unlimit / VM', 'limit / VM', 'p99 unlimit Δ', 'p99 limit Δ'],
          rows=[
            [[('TiDB v8.5.2', {'bold': True})],
             '26,947', '23,442.9', '15,751.9',
             [('87.0 %', {'bold': True, 'color': RED_EMPH})],
             '58.5 %', '+17 %', '+111 %'],
            [[('CockroachDB v26.2', {'bold': True})],
             '15,033', '12,196.7', '6,493.5',
             [('81.1 %', {'bold': True, 'color': RED_EMPH})],
             '43.2 %', '+27 %', '+192 %'],
            [[('YugabyteDB 2025.2', {'bold': True})],
             '15,632', '2,997.6', '1,604.5',
             [('19.2 % ⚠', {'bold': True, 'color': RED_EMPH})],
             '10.3 % ⚠', '+669 %', '+1556 %'],
          ],
          col_widths=[2.2, 1.5, 1.5, 1.5, 1.4, 1.2, 1.4, 1.4],
          header_size=12, body_size=13)
add_text(s, Inches(0.6), Inches(4.85), Inches(12), Inches(0.5),
         '觀察解讀', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(0.8), Inches(5.3), Inches(12), Inches(1.6),
            ['TiDB / CockroachDB：K8s unlimit 保留率 81-87 %，屬可接受範圍；limit 情境大幅衰退反映資源壓縮',
             [('YugabyteDB：K8s unlimit 僅 19.2 %、p99 +669 %。成因推論：YSQL + DocDB 雙 process 在 K8s pod IPC + CPU contention 放大。',
               {}),
              ('成因未定位，列後續調校項', {'bold': True, 'color': RED_EMPH})],
             'S-BASE ↔ S-K8S 可直接比較（同 workload W=128 t=128）'],
            size=15, color=SLATE)

# =========================================================================
# Slide 8 — 跨區域 / 跨專線進度
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 8, 11)
add_title(s, '跨區域 / 跨專線進度')
add_runs(s, Inches(0.75), Inches(1.4), Inches(12), Inches(0.4),
         runs=[('{FIXME}', {'bold': True, 'color': RED_EMPH, 'size': 14}),
               (' 架構圖示說明解釋 P-A/P-B placement', {'color': MUTED, 'size': 14})])
# Bar before subtitle
add_rect(s, Inches(0.6), Inches(1.4), Pt(3), Inches(0.35), ACCENT_RED)

add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
         '✅ 已完成', size=20, bold=True, color=RGBColor(0x27, 0xAE, 0x60))
add_bullets(s, Inches(0.9), Inches(2.5), Inches(12), Inches(2.2),
            ['技術議題定案（GCP 5 VM 拓樸、placement P-A/P-B、Chaos Engineering 等）',
             'IaC、playbook、suite scripts、chrony drift gate（drift_median 0.017 ms；時間偏移 between IDC & GCP）',
             '06-18〜06-19：修正 IaC / 防火牆 / YugabyteDB placement，三家 DB / 六節點前期測試驗證',
             '06-21〜06-22：Determinism v2 — 同 cluster R1-R5、scheduler / balancer freeze'],
            size=16, color=SLATE)

add_text(s, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
         '● 待執行', size=20, bold=True, color=yellow)
add_bullets(s, Inches(0.9), Inches(5.3), Inches(12), Inches(1.6),
            ['調度 scheduler (開/關)、round-only runner、Warmup、Placement (replica / lease 位置)、測試結果彙整',
             '回到 W=128 正式基準量測',
             'Failover 與 Chaos Engineering 實際驗證'],
            size=16, color=SLATE)

# =========================================================================
# Slide 9 — X-CROSS 初步結果
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 9, 11)
add_title(s, 'X-CROSS 初步結果 (06-22〜06-23)')
add_subtitle(s, 'W=4 framework 驗證；不與 S-BASE/S-K8S 直接對比 · 同 cluster 5 rounds、W=4、16 threads、每 round 5 min、controller = .31 (IDC client)',
             Inches(1.4), color=MUTED, size=13)
green_check = RGBColor(0x27, 0xAE, 0x60)
add_table(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(2.1),
          header=['資料庫', 'R1', 'R2', 'R3', 'R4', 'R5', '中位數 (有效輪)', 'CV'],
          rows=[
            [[('TiDB v8.5.2', {'bold': True})],
             '9,525.5', '9,553.2', '9,786.9', '9,393.2', '9,530.8',
             [('9,566.0 (R2-R5)', {'bold': True, 'color': RED_EMPH})],
             [('1.67 % ✅', {'bold': True, 'color': green_check})]],
            [[('CockroachDB v26.2', {'bold': True})],
             '8,409.5', '8,055.3', '7,902.5', '7,720.9', '7,472.3',
             [('7,787.8 (R2-R5)', {'bold': True, 'color': RED_EMPH})],
             [('3.23 % ✅', {'bold': True, 'color': green_check})]],
            [[('YugabyteDB 2025.2', {'bold': True})],
             '102.0', '226.9', '6,424.2', '6,259.3', '6,206.2',
             [('6,296.6 (R3-R5)', {'bold': True, 'color': RED_EMPH})],
             [('1.82 % ✅', {'bold': True, 'color': green_check})]],
          ],
          col_widths=[2.4, 1.1, 1.1, 1.1, 1.1, 1.1, 2.3, 1.3],
          header_size=12, body_size=12)
add_text(s, Inches(0.6), Inches(4.35), Inches(12), Inches(0.4),
         '已確認結論', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(0.8), Inches(4.8), Inches(12.3), Inches(1.7),
            ['06-21 觀察到的 ±526 % / ±50 % 變異主因為「每輪重新部署」造成 placement / cache / scheduler 狀態異動，與 W=4 contention 為兩個獨立來源',
             '同一 cluster 連跑時，三家 W=4 R2-R5（或 YBDB R3-R5）的 CV ≤ 5 %，重現性已建立',
             'YBDB Idle=0 解法：gate 改用 get_load_move_completion=100%；timed run 前 set_load_balancer_enabled=0',
             'CRDB lease gate SQL 改用 SHOW RANGES FROM DATABASE tpcc WITH TABLES, DETAILS（v26.2 相容）'],
            size=13, color=SLATE)
add_runs(s, Inches(0.6), Inches(6.6), Inches(12.5), Inches(0.4),
         runs=[('⚠ Caveat：', {'bold': True, 'color': yellow, 'size': 13}),
               ('W=4 framework / contention 驗證 ≠ 正式基準；R2-R5（YBDB R3-R5）CV ≤ 5 % 已建立但',
                {'color': MUTED, 'size': 13}),
               ('不可作跨家排名', {'bold': True, 'color': RED_EMPH, 'size': 13}),
               ('；後續需回 W=128。', {'color': MUTED, 'size': 13})])

# =========================================================================
# Slide 10 — 建議決策框架
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 10, 11)
add_title(s, '建議決策框架')
add_subtitle(s, '不是選一家，是先界定 application 條件、再排序候選；僅提供做決策框架參考依據用',
             Inches(1.4), color=MUTED)
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.0),
          header=['分類', '內容', '處置'],
          rows=[
            [[('短期候選', {'bold': True, 'color': RGBColor(0x27, 0xAE, 0x60)})],
             [('TiDB', {'bold': True})],
             '優先進入應用情境對接，無門檻'],
            [[('保留觀察', {'bold': True, 'color': ORANGE})],
             'CockroachDB',
             '依一致性需求、Error Handling 容忍度、維運成本評估'],
            [[('保留觀察', {'bold': True, 'color': ORANGE})],
             'YugabyteDB',
             'VM 路徑可評估；K8s 路徑需先完成部署層級調校與驗證'],
            [[('暫不作結論', {'bold': True, 'color': RED_EMPH})],
             '跨區域場景、K8s 退化未定位項',
             '待下一階段 W=128 測試數據與調校產出再回頭評估'],
          ],
          col_widths=[1.8, 3.5, 6.8])
add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
         'Application 需要確認的議題', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(0.8), Inches(4.8), Inches(12.3), Inches(2.3),
            [[('1. 交易一致性需求：', {'bold': True, 'color': ORANGE}),
              ('是否需要 SERIALIZABLE / SSI？是否可接受 READ COMMITTED & 產品設計架構轉 CAP 的複雜度及可行性？', {})],
             [('2. 可接受延遲：', {'bold': True, 'color': ORANGE}),
              ('在尖峰 t=128 等級負載下，p99 是否需要 500ms / 1s / 2s 以內？', {})],
             [('3. retry / timeout 行為：', {'bold': True, 'color': ORANGE}),
              ('應用層是否能接受 SI / SSI 模式下的 retry 機制？', {})],
             [('4. RTO / RPO：', {'bold': True, 'color': ORANGE}),
              ('跨區域 failover 是否需要 < 30s RTO？是否容許資料 lag？', {})],
             [('5. 連線層 / 交易模式調整：', {'bold': True, 'color': ORANGE}),
              ('是否能配合 HAProxy / PgBouncer / 連線池與短交易模式？', {})]],
            size=13, color=SLATE)

# =========================================================================
# Slide 11 — 後續推進階段
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_NAVY)
add_header_footer(s, 11, 11)
add_title(s, '後續推進階段')
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.1),
          header=['短期 EX: Y26/07', '中期', '決策'],
          rows=[
            [[('{FIXME}', {'bold': True, 'color': RED_EMPH, 'size': 13}),
              (' 補完完整測試時間描述', {'size': 13})],
             'A-S / A-A-RO / A-A 三個 workload profile × P-A/P-B placement × 三家資料庫，依檢驗指標分批 / 量測',
             'Application owner 完成五項議題確認'],
            ['W=4 框架已驗，接 W=128 基準',
             'Failover 與 Chaos Engineering 測試',
             '原廠後勤對接狀況說明'],
            ['freeze/unfreeze、round-only runner、暖機、placement gate 全部通過 審查',
             '跨區域 analytics 第二份報告',
             ''],
            ['X-CROSS 完成 determinism 流程 驗收，啟動正式 W=128 基準 量測',
             '補齊 YugabyteDB K8s 退化成因調查（P-A/P-B、A-S/A-A-RO/A-A）',
             ''],
          ],
          col_widths=[4.5, 5.0, 2.6],
          header_size=14, body_size=13)

# =========================================================================
# Save
# =========================================================================
out = '/Users/wn.lin/vscode-git/dba_career/poc/1_MeetingMinutes/0626-slide-v5.pptx'
prs.save(out)
print(f'wrote {out}')
